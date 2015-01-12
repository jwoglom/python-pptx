# encoding: utf-8

"""
Test suite for pptx.shapes.placeholder module
"""

from __future__ import absolute_import, print_function, unicode_literals

import pytest

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.shapes.shared import ST_Direction, ST_PlaceholderSize
from pptx.parts.image import ImagePart
from pptx.parts.slide import Slide
from pptx.parts.slidelayout import SlideLayout
from pptx.parts.slidemaster import SlideMaster
from pptx.shapes.placeholder import (
    BasePlaceholder, _BaseSlidePlaceholder, LayoutPlaceholder,
    MasterPlaceholder, PicturePlaceholder, PlaceholderGraphicFrame,
    PlaceholderPicture, TablePlaceholder
)
from pptx.shapes.shapetree import BaseShapeTree

from ..oxml.unitdata.shape import (
    an_ext, a_graphicFrame, a_ph, an_nvGraphicFramePr, an_nvPicPr, an_nvPr,
    an_nvSpPr, an_sp, an_spPr, an_xfrm
)
from ..unitutil.cxml import element, xml
from ..unitutil.mock import (
    class_mock, instance_mock, method_mock, property_mock
)


class Describe_BaseSlidePlaceholder(object):

    def it_knows_its_shape_type(self):
        placeholder = _BaseSlidePlaceholder(None, None)
        assert placeholder.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER

    def it_provides_override_dimensions_when_present(self, override_fixture):
        placeholder, prop_name, expected_value = override_fixture
        assert getattr(placeholder, prop_name) == expected_value

    def it_provides_inherited_dims_when_no_override(self, inherited_fixture):
        placeholder, prop_name, expected_value = inherited_fixture
        value = getattr(placeholder, prop_name)
        placeholder._inherited_value.assert_called_once_with(prop_name)
        assert value == expected_value

    def it_gets_an_inherited_dim_value_to_help(self, layout_val_fixture):
        placeholder, attr_name, expected_value = layout_val_fixture
        value = placeholder._inherited_value(attr_name)
        assert value == expected_value

    def it_finds_its_layout_placeholder_to_help(self, layout_ph_fixture):
        placeholder, layout_, idx, layout_placeholder_ = layout_ph_fixture
        layout_placeholder = placeholder._layout_placeholder
        layout_.placeholders.get.assert_called_once_with(idx=idx)
        assert layout_placeholder is layout_placeholder_

    def it_finds_its_slide_layout_to_help(self, slide_layout_fixture):
        placeholder, slide_layout_ = slide_layout_fixture
        slide_layout = placeholder._slide_layout
        assert slide_layout is slide_layout_

    def it_can_override_inherited_dimensions(self, dim_set_fixture):
        placeholder, prop_name, value, expected_xml = dim_set_fixture
        setattr(placeholder, prop_name, value)
        assert placeholder._element.xml == expected_xml

    def it_replaces_a_placeholder_element_to_help(self, replace_fixture):
        placeholder, element, spTree, expected_xml = replace_fixture
        placeholder._replace_placeholder_with(element)
        assert spTree.xml == expected_xml
        assert placeholder._element is None

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sp/p:spPr/a:xfrm', 'left',   1,
         'p:sp/p:spPr/a:xfrm/a:off{x=1,y=0}'),
        ('p:sp/p:spPr/a:xfrm', 'top',    2,
         'p:sp/p:spPr/a:xfrm/a:off{x=0,y=2}'),
        ('p:sp/p:spPr/a:xfrm', 'width',  3,
         'p:sp/p:spPr/a:xfrm/a:ext{cx=3,cy=0}'),
        ('p:sp/p:spPr/a:xfrm', 'height', 4,
         'p:sp/p:spPr/a:xfrm/a:ext{cx=0,cy=4}'),
    ])
    def dim_set_fixture(self, request):
        sp_cxml, prop_name, value, expected_cxml = request.param
        placeholder = _BaseSlidePlaceholder(element(sp_cxml), None)
        expected_xml = xml(expected_cxml)
        return placeholder, prop_name, value, expected_xml

    @pytest.fixture(params=['left', 'top', 'width', 'height'])
    def inherited_fixture(self, request, _inherited_value_):
        prop_name = request.param
        placeholder = _BaseSlidePlaceholder(element('p:sp/p:spPr'), None)
        _inherited_value_.return_value = expected_value = 42
        return placeholder, prop_name, expected_value

    @pytest.fixture
    def layout_ph_fixture(self, request, _slide_layout_prop_, slide_layout_,
                          layout_placeholder_):
        sp_cxml = 'p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=1}'
        placeholder = _BaseSlidePlaceholder(element(sp_cxml), None)
        _slide_layout_prop_.return_value = slide_layout_
        slide_layout_.placeholders.get.return_value = layout_placeholder_
        return placeholder, slide_layout_, 1, layout_placeholder_

    @pytest.fixture(params=[(True, 42), (False, None)])
    def layout_val_fixture(self, request, _layout_placeholder_,
                           layout_placeholder_):
        has_layout_placeholder, expected_value = request.param
        placeholder = _BaseSlidePlaceholder(None, None)
        if has_layout_placeholder:
            layout_placeholder_.width = expected_value
            _layout_placeholder_.return_value = layout_placeholder_
        else:
            _layout_placeholder_.return_value = None
        return placeholder, 'width', expected_value

    @pytest.fixture(params=[
        ('left',   'p:sp/p:spPr/a:xfrm/a:off{x=12}',  12),
        ('top',    'p:sp/p:spPr/a:xfrm/a:off{y=34}',  34),
        ('width',  'p:sp/p:spPr/a:xfrm/a:ext{cx=56}', 56),
        ('height', 'p:sp/p:spPr/a:xfrm/a:ext{cy=78}', 78),
    ])
    def override_fixture(self, request):
        prop_name, sp_cxml, value = request.param
        slide_placeholder = _BaseSlidePlaceholder(element(sp_cxml), None)
        return slide_placeholder, prop_name, value

    @pytest.fixture
    def replace_fixture(self):
        spTree = element(
            'p:spTree/p:sp/p:nvSpPr/p:nvPr/p:ph{type=pic,idx=10}'
        )
        pic = element('p:pic/p:nvPicPr/p:nvPr')
        placeholder = _BaseSlidePlaceholder(spTree[0], None)
        expected_xml = xml(
            'p:spTree/p:pic/p:nvPicPr/p:nvPr/p:ph{type=pic,idx=10}'
        )
        return placeholder, pic, spTree, expected_xml

    @pytest.fixture
    def slide_layout_fixture(self, part_prop_, slide_, slide_layout_):
        placeholder = _BaseSlidePlaceholder(None, None)
        part_prop_.return_value = slide_
        slide_.slide_layout = slide_layout_
        return placeholder, slide_layout_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _inherited_value_(self, request):
        return method_mock(
            request, _BaseSlidePlaceholder, '_inherited_value',
        )

    @pytest.fixture
    def _layout_placeholder_(self, request):
        return property_mock(
            request, _BaseSlidePlaceholder, '_layout_placeholder'
        )

    @pytest.fixture
    def layout_placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, _BaseSlidePlaceholder, 'part')

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def _slide_layout_prop_(self, request):
        return property_mock(request, _BaseSlidePlaceholder, '_slide_layout')


class DescribeBasePlaceholder(object):

    def it_knows_its_idx_value(self, idx_fixture):
        placeholder, idx = idx_fixture
        assert placeholder.idx == idx

    def it_knows_its_orient_value(self, orient_fixture):
        placeholder, orient = orient_fixture
        assert placeholder.orient == orient

    def it_raises_on_ph_orient_when_not_a_placeholder(
            self, orient_raise_fixture):
        shape = orient_raise_fixture
        with pytest.raises(ValueError):
            shape.orient

    def it_knows_its_sz_value(self, sz_fixture):
        placeholder, sz = sz_fixture
        assert placeholder.sz == sz

    def it_knows_its_placeholder_type(self, type_fixture):
        placeholder, ph_type = type_fixture
        assert placeholder.ph_type == ph_type

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('sp',           None, 'title', 0),
        ('sp',           0,    'title', 0),
        ('sp',           3,    'body',  3),
        ('pic',          6,    'pic',   6),
        ('graphicFrame', 9,    'tbl',   9),
    ])
    def idx_fixture(self, request):
        tagname, idx, ph_type, expected_idx = request.param
        shape_elm = self.shape_elm_factory(tagname, ph_type, idx)
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_idx

    @pytest.fixture(params=[
        (None, ST_Direction.HORZ),
        (ST_Direction.VERT, ST_Direction.VERT),
    ])
    def orient_fixture(self, request):
        orient, expected_orient = request.param
        ph_bldr = a_ph()
        if orient is not None:
            ph_bldr.with_orient(orient)
        shape_elm = (
            an_sp().with_nsdecls('p').with_child(
                an_nvSpPr().with_child(
                    an_nvPr().with_child(
                        ph_bldr)))
        ).element
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_orient

    @pytest.fixture
    def orient_raise_fixture(self):
        shape_elm = (
            an_sp().with_nsdecls('p').with_child(
                an_nvSpPr().with_child(
                    an_nvPr()))
        ).element
        shape = BasePlaceholder(shape_elm, None)
        return shape

    @pytest.fixture(params=[
        (None, ST_PlaceholderSize.FULL),
        (ST_PlaceholderSize.HALF, ST_PlaceholderSize.HALF),
    ])
    def sz_fixture(self, request):
        sz, expected_sz = request.param
        ph_bldr = a_ph()
        if sz is not None:
            ph_bldr.with_sz(sz)
        shape_elm = (
            an_sp().with_nsdecls('p').with_child(
                an_nvSpPr().with_child(
                    an_nvPr().with_child(
                        ph_bldr)))
        ).element
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_sz

    @pytest.fixture(params=[
        ('sp',           None,    1, PP_PLACEHOLDER.OBJECT),
        ('sp',           'title', 0, PP_PLACEHOLDER.TITLE),
        ('pic',          'pic',   6, PP_PLACEHOLDER.PICTURE),
        ('graphicFrame', 'tbl',   9, PP_PLACEHOLDER.TABLE),
    ])
    def type_fixture(self, request):
        tagname, ph_type, idx, expected_ph_type = request.param
        shape_elm = self.shape_elm_factory(tagname, ph_type, idx)
        placeholder = BasePlaceholder(shape_elm, None)
        return placeholder, expected_ph_type

    # fixture components ---------------------------------------------

    @staticmethod
    def shape_elm_factory(tagname, ph_type, idx):
        root_bldr, nvXxPr_bldr = {
            'sp':           (an_sp().with_nsdecls('p'), an_nvSpPr()),
            'pic':          (an_sp().with_nsdecls('p'), an_nvPicPr()),
            'graphicFrame': (a_graphicFrame().with_nsdecls('p'),
                             an_nvGraphicFramePr()),
        }[tagname]
        ph_bldr = {
            None:    a_ph().with_idx(idx),
            'obj':   a_ph().with_idx(idx),
            'title': a_ph().with_type('title'),
            'body':  a_ph().with_type('body').with_idx(idx),
            'pic':   a_ph().with_type('pic').with_idx(idx),
            'tbl':   a_ph().with_type('tbl').with_idx(idx),
        }[ph_type]
        return (
            root_bldr.with_child(
                nvXxPr_bldr.with_child(
                    an_nvPr().with_child(
                        ph_bldr)))
        ).element


class DescribeLayoutPlaceholder(object):

    def it_considers_inheritance_when_computing_pos_and_size(
            self, xfrm_fixture):
        layout_placeholder, _direct_or_inherited_value_ = xfrm_fixture[:2]
        attr_name, expected_value = xfrm_fixture[2:]
        value = getattr(layout_placeholder, attr_name)
        _direct_or_inherited_value_.assert_called_once_with(attr_name)
        assert value == expected_value

    def it_provides_direct_property_values_when_they_exist(
            self, direct_fixture):
        layout_placeholder, expected_width = direct_fixture
        width = layout_placeholder.width
        assert width == expected_width

    def it_provides_inherited_property_values_when_no_direct_value(
            self, inherited_fixture):
        layout_placeholder, _inherited_value_, inherited_left_ = (
            inherited_fixture
        )
        left = layout_placeholder.left
        _inherited_value_.assert_called_once_with('left')
        assert left == inherited_left_

    def it_knows_how_to_get_a_property_value_from_its_master(
            self, mstr_val_fixture):
        layout_placeholder, attr_name, expected_value = mstr_val_fixture
        value = layout_placeholder._inherited_value(attr_name)
        assert value == expected_value

    def it_finds_its_corresponding_master_placeholder_to_help_inherit(
            self, mstr_ph_fixture):
        layout_placeholder, master_, mstr_ph_type, master_placeholder_ = (
            mstr_ph_fixture
        )
        master_placeholder = layout_placeholder._master_placeholder
        master_.placeholders.get.assert_called_once_with(mstr_ph_type, None)
        assert master_placeholder is master_placeholder_

    def it_finds_its_slide_master_to_help_inherit(self, slide_master_fixture):
        layout_placeholder, slide_master_ = slide_master_fixture
        slide_master = layout_placeholder._slide_master
        assert slide_master == slide_master_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def direct_fixture(self, sp, width):
        layout_placeholder = LayoutPlaceholder(sp, None)
        return layout_placeholder, width

    @pytest.fixture
    def inherited_fixture(self, sp, _inherited_value_, int_value_):
        layout_placeholder = LayoutPlaceholder(sp, None)
        return layout_placeholder, _inherited_value_, int_value_

    @pytest.fixture(params=[
        (PP_PLACEHOLDER.TABLE,  PP_PLACEHOLDER.BODY),
        (PP_PLACEHOLDER.BODY,   PP_PLACEHOLDER.BODY),
        (PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY),
    ])
    def mstr_ph_fixture(
            self, request, ph_type_, _slide_master_, slide_master_,
            master_placeholder_):
        layout_placeholder = LayoutPlaceholder(None, None)
        ph_type, mstr_ph_type = request.param
        ph_type_.return_value = ph_type
        slide_master_.placeholders.get.return_value = master_placeholder_
        return (
            layout_placeholder, slide_master_, mstr_ph_type,
            master_placeholder_
        )

    @pytest.fixture(params=[(True, 42), (False, None)])
    def mstr_val_fixture(
            self, request, _master_placeholder_, master_placeholder_):
        has_master_placeholder, expected_value = request.param
        layout_placeholder = LayoutPlaceholder(None, None)
        attr_name = 'width'
        if has_master_placeholder:
            setattr(master_placeholder_, attr_name, expected_value)
            _master_placeholder_.return_value = master_placeholder_
        else:
            _master_placeholder_.return_value = None
        return layout_placeholder, attr_name, expected_value

    @pytest.fixture
    def slide_master_fixture(self, parent_, slide_master_):
        layout_placeholder = LayoutPlaceholder(None, parent_)
        return layout_placeholder, slide_master_

    @pytest.fixture(params=['left', 'top', 'width', 'height'])
    def xfrm_fixture(self, request, _direct_or_inherited_value_, int_value_):
        attr_name = request.param
        layout_placeholder = LayoutPlaceholder(None, None)
        _direct_or_inherited_value_.return_value = int_value_
        return (
            layout_placeholder, _direct_or_inherited_value_, attr_name,
            int_value_
        )

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _direct_or_inherited_value_(self, request):
        return method_mock(
            request, LayoutPlaceholder, '_direct_or_inherited_value'
        )

    @pytest.fixture
    def _inherited_value_(self, request, int_value_):
        return method_mock(
            request, LayoutPlaceholder, '_inherited_value',
            return_value=int_value_
        )

    @pytest.fixture
    def int_value_(self, request):
        return instance_mock(request, int)

    @pytest.fixture
    def _master_placeholder_(self, request):
        return property_mock(
            request, LayoutPlaceholder, '_master_placeholder'
        )

    @pytest.fixture
    def master_placeholder_(self, request):
        return instance_mock(request, MasterPlaceholder)

    @pytest.fixture
    def parent_(self, request, slide_layout_):
        parent_ = instance_mock(request, BaseShapeTree)
        parent_.part = slide_layout_
        return parent_

    @pytest.fixture
    def ph_type_(self, request):
        return property_mock(request, LayoutPlaceholder, 'ph_type')

    @pytest.fixture
    def slide_layout_(self, request, slide_master_):
        slide_layout_ = instance_mock(request, SlideLayout)
        slide_layout_.slide_master = slide_master_
        return slide_layout_

    @pytest.fixture
    def _slide_master_(self, request, slide_master_):
        return property_mock(
            request, LayoutPlaceholder, '_slide_master',
            return_value=slide_master_
        )

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)

    @pytest.fixture
    def sp(self, width):
        return (
            an_sp().with_nsdecls('p', 'a').with_child(
                an_spPr().with_child(
                    an_xfrm().with_child(
                        an_ext().with_cx(width))))
        ).element

    @pytest.fixture
    def width(self):
        return 31416


class DescribePicturePlaceholder(object):

    def it_can_insert_a_picture_into_itself(self, insert_fixture):
        picture_ph, image_file, pic = insert_fixture[:3]
        PlaceholderPicture_, placeholder_picture_ = insert_fixture[3:]

        placeholder_picture = picture_ph.insert_picture(image_file)

        picture_ph._new_placeholder_pic.assert_called_once_with(image_file)
        picture_ph._replace_placeholder_with.assert_called_once_with(pic)
        PlaceholderPicture_.assert_called_once_with(pic, picture_ph._parent)
        assert placeholder_picture is placeholder_picture_

    def it_creates_a_pic_element_to_help(self, pic_fixture):
        picture_ph, image_file, expected_xml = pic_fixture
        pic = picture_ph._new_placeholder_pic(image_file)
        picture_ph._get_or_add_image.assert_called_once_with(image_file)
        assert pic.xml == expected_xml

    def it_adds_an_image_to_help(self, get_or_add_fixture):
        placeholder, image_file, expected_value = get_or_add_fixture

        value = placeholder._get_or_add_image(image_file)

        placeholder.part.get_or_add_image_part.assert_called_once_with(
            image_file
        )
        assert value == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def get_or_add_fixture(self, part_prop_, image_part_):
        placeholder = PicturePlaceholder(None, None)
        image_file, rId, desc, image_size = 'f.png', 'rId6', 'desc', (42, 24)
        part_prop_.return_value.get_or_add_image_part.return_value = (
            image_part_, rId
        )
        image_part_.desc, image_part_._px_size = desc, image_size
        expected_value = rId, desc, image_size
        return placeholder, image_file, expected_value

    @pytest.fixture
    def insert_fixture(self, PlaceholderPicture_, placeholder_picture_,
                       _new_placeholder_pic_, _replace_placeholder_with_):
        picture_ph = PicturePlaceholder(None, 'parent')
        image_file, pic = 'foobar.png', element('p:pic')
        _new_placeholder_pic_.return_value = pic
        PlaceholderPicture_.return_value = placeholder_picture_
        return (
            picture_ph, image_file, pic, PlaceholderPicture_,
            placeholder_picture_
        )

    @pytest.fixture(params=[
        ((444, 333), ('l', 'r')),
        ((333, 444), ('t', 'b')),
    ])
    def pic_fixture(self, request, _get_or_add_image_):
        image_size, crop_attr_names = request.param
        sp_cxml = (
            'p:sp/(p:nvSpPr/p:cNvPr{id=2,name=foo},p:spPr/a:xfrm/a:ext{cx=99'
            ',cy=99})'
        )
        sp = element(sp_cxml)
        picture_ph = PicturePlaceholder(sp, None)
        image_file = 'foobar.png'
        _get_or_add_image_.return_value = 42, 'bar', image_size
        expected_xml = xml(
            'p:pic/(p:nvPicPr/(p:cNvPr{id=2,name=foo,descr=bar},p:cNvPicPr/a'
            ':picLocks{noGrp=1,noChangeAspect=1},p:nvPr),p:blipFill/(a:blip{'
            'r:embed=42},a:srcRect{%s=12500,%s=12500},a:stretch/a:fillRect),'
            'p:spPr)' % crop_attr_names
        )
        return picture_ph, image_file, expected_xml

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _get_or_add_image_(self, request):
        return method_mock(
            request, PicturePlaceholder, '_get_or_add_image'
        )

    @pytest.fixture
    def image_part_(self, request):
        return instance_mock(request, ImagePart)

    @pytest.fixture
    def _new_placeholder_pic_(self, request):
        return method_mock(
            request, PicturePlaceholder, '_new_placeholder_pic'
        )

    @pytest.fixture
    def part_prop_(self, request, slide_):
        return property_mock(
            request, PicturePlaceholder, 'part', return_value=slide_
        )

    @pytest.fixture
    def PlaceholderPicture_(self, request):
        return class_mock(
            request, 'pptx.shapes.placeholder.PlaceholderPicture'
        )

    @pytest.fixture
    def placeholder_picture_(self, request):
        return instance_mock(request, PlaceholderPicture)

    @pytest.fixture
    def _replace_placeholder_with_(self, request):
        return method_mock(
            request, PicturePlaceholder, '_replace_placeholder_with'
        )

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)


class DescribeTablePlaceholder(object):

    def it_can_insert_a_table_into_itself(self, insert_fixture):
        table_ph, rows, cols, graphicFrame = insert_fixture[:4]
        PlaceholderGraphicFrame_, ph_graphic_frame_ = insert_fixture[4:]

        ph_graphic_frame = table_ph.insert_table(rows, cols)

        table_ph._new_placeholder_table.assert_called_once_with(rows, cols)
        table_ph._replace_placeholder_with.assert_called_once_with(
            graphicFrame
        )
        PlaceholderGraphicFrame_.assert_called_once_with(
            graphicFrame, table_ph._parent
        )
        assert ph_graphic_frame is ph_graphic_frame_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def insert_fixture(
            self, PlaceholderGraphicFrame_, placeholder_graphic_frame_,
            _new_placeholder_table_, _replace_placeholder_with_):
        table_ph = TablePlaceholder(None, 'parent')
        rows, cols, graphicFrame = 4, 2, element('p:graphicFrame')
        _new_placeholder_table_.return_value = graphicFrame
        PlaceholderGraphicFrame_.return_value = placeholder_graphic_frame_
        return (
            table_ph, rows, cols, graphicFrame, PlaceholderGraphicFrame_,
            placeholder_graphic_frame_
        )

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _new_placeholder_table_(self, request):
        return method_mock(
            request, TablePlaceholder, '_new_placeholder_table'
        )

    @pytest.fixture
    def PlaceholderGraphicFrame_(self, request):
        return class_mock(
            request, 'pptx.shapes.placeholder.PlaceholderGraphicFrame'
        )

    @pytest.fixture
    def placeholder_graphic_frame_(self, request):
        return instance_mock(request, PlaceholderGraphicFrame)

    @pytest.fixture
    def _replace_placeholder_with_(self, request):
        return method_mock(
            request, TablePlaceholder, '_replace_placeholder_with'
        )
